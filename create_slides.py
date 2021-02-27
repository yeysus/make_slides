# Create a set of PowerPoint slides based on Screenshots from Websites and
# Info text in an Excel file.

# sudo pip3 install Pillow
# sudo pip3 install lxml
# sudo pip3 install XlsxWriter
# sudo pip3 install python-pptx

# How to place pictures in a slide:
# https://stackoverflow.com/questions/44275443/python-inserts-pictures-to-powerpoint-how-to-set-the-width-and-height-of-the-pi

# How to calculate aspect ratio:
# https://www.omnicalculator.com/other/aspect-ratio

import sys
# For sleep() in some tests I did.
import time
import pandas as pd
import os, sys, subprocess
import json

from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor
from pptx.util import Inches, Pt

def open_file(filename):
  if sys.platform == "win32":
    os.startfile(filename)
  else:
    opener = "open" if sys.platform == "darwin" else "xdg-open"
    subprocess.call([opener, filename])

# From: https://stackoverflow.com/questions/287871/how-to-print-colored-text-to$
class bcolors:
  HEADER = '\033[95m'
  OKBLUE = '\033[94m'
  OKCYAN = '\033[96m'
  OKGREEN = '\033[92m'
  WARNING = '\033[93m'
  FAIL = '\033[91m'
  ENDC = '\033[0m'
  BOLD = '\033[1m'
  UNDERLINE = '\033[4m'

screenshotsDirectory = '/Users/jesusdelvalle/Documents/projects/make_slides/screenshots/'

nameOfExcelColumnWithURLs = 'URL'
nameOfExcelColumnWithCompanyName = 'Name'
nameOfExcelColumnWithYearFounded = 'Founded'
nameOfExcelColumnWithHQCity = 'HQ'
nameOfExcelColumnWithDescription = 'Description'

# This title will appear on the top of every slide for certain formats.
title = "Top30 BRITISH / WHERE to Make an Internship 2021"

powerpoint = Presentation ("examples/Template_16x9.pptx")
finalPowerPoint = "examples/Slides.pptx"

# The following measures are in inches.
# titleHeight is the height reserved for the title, if any.
# The title font size could be smaller.
# The total size is 16 x 9.

# Example Layout: layout 1.
# Title
# 2 Pictures per slide, both left, with corresponding text on their right.
#####
# Title
#--------
#| pic1 | text1
#--------
#| pic2 | text2
#--------
layout1 = '{"layout":"1", \
           "referenceHeight":3, \
           "totalHeight":9, \
           "nColumnsPerSlide":1, \
           "textBoxPosition":"right", \
           "pictureNameFont":40, \
           "pictureDescriptionFont":20, \
           "pictureUrlFont":12, \
           "titleHeight":1}'

layout2 = '{"layout":"2", \
           "referenceHeight":3.65, \
           "totalHeight":9, \
           "nColumnsPerSlide":2, \
           "textBoxPosition":"overlap", \
           "pictureNameFont":28, \
           "pictureDescriptionFont":20, \
           "pictureUrlFont":12, \
           "titleHeight":0}'

layoutParameters = json.loads(layout2)

totalHeight = layoutParameters ["totalHeight"]
titleHeight = layoutParameters ["titleHeight"]
referenceHeight = layoutParameters ["referenceHeight"]
nColumnsPerSlide = layoutParameters ["nColumnsPerSlide"]
textBoxPosition = layoutParameters ["textBoxPosition"]
pictureNameFont = layoutParameters ["pictureNameFont"]
pictureDescriptionFont = layoutParameters ["pictureDescriptionFont"]
pictureUrlFont = layoutParameters ["pictureUrlFont"]

pictureHeight = Inches (referenceHeight)
pictureWidth = Inches (referenceHeight*9/5)
# Total height is 9 inches. Reserving titleHeight for the title,
# nRowsPerSlide = (9 - titleHeight) // referenceHeight
# In Python, the ( // ) operator returns the integer quotient.
nRowsPerSlide = (totalHeight - titleHeight) // referenceHeight

try:
  # Expected: Excel file with ending .xslx
  inputFilename = sys.argv[1]

  excel = pd.read_excel(inputFilename, header=0)

  # i is the counter for the number of items which will go into the slides.
  i = 0
  # iRow is the counter for the number of rows which will go into each slide.
  iRow = 0
  # iColumn is the counter for the number of columns which will go into each slide.
  iColumn = 0
  layout = powerpoint.slide_layouts[0]
  slide = powerpoint.slides.add_slide(layout)
  shape = slide.shapes
  for index, row in excel.iterrows():

    url = row [nameOfExcelColumnWithURLs]

    url = url.replace("https://","")
    url = url.replace("/","")
    print (url)
    pictureFile = screenshotsDirectory + url + ".png"

    if iRow == 0 & iColumn == 0:
      slide = powerpoint.slides.add_slide(layout)
      if titleHeight > 0:
        textBox = slide.shapes.add_textbox (Inches (0.05), Inches (0.0), Inches(14), Inches(1))
        textFrame = textBox.text_frame
        textFrame.word_wrap = True
        # A text_frame always contains one paragraph.
        p = textFrame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(245, 223, 77)
        p.text = title

    # 16:9 -> 3 columns each 5.33, 4 rows each 2.25
    # new height = (5/9) * new width

    # 13.3 x 7.5 inches -> 4.43 x 1.875

    left = Inches (iColumn * pictureWidth + (iColumn / 4.5))
    top = Inches (titleHeight + (iRow * referenceHeight) + (iRow / 4.5))

    picture=slide.shapes.add_picture (pictureFile, left, top, width = pictureWidth, height = pictureHeight)

    # Add text.
    # From https://python-pptx.readthedocs.io/en/latest/user/text.html

    if textBoxPosition == "right":
      left = left + pictureWidth
    elif textBoxPosition == "overlap":
      left = left
    txBox = slide.shapes.add_textbox (left, top, Inches(totalHeight - titleHeight), Inches(referenceHeight))
    tf = txBox.text_frame
    tf.word_wrap = True
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor (0, 0, 0)

    p = tf.add_paragraph()
    p.font.size = Pt(pictureNameFont)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.text = row [nameOfExcelColumnWithCompanyName]
    p = tf.add_paragraph()

    # Clean missing values.
    rowFounded = ""
    if pd.isnull (row [nameOfExcelColumnWithYearFounded]):
      rowFounded = ""
    else:
      rowFounded = "Founded: " + str (int (row [nameOfExcelColumnWithYearFounded])) + ". "
    p.font.size = Pt(pictureDescriptionFont)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.text = row [nameOfExcelColumnWithHQCity] + '. ' + rowFounded + row [nameOfExcelColumnWithDescription]

    p = tf.add_paragraph()
    p.font.size = Pt(pictureUrlFont)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.text = row [nameOfExcelColumnWithURLs]

    #j = j + 1
    # print ("Shape written: " + str (j))
    iRow = iRow + 1
    # print ("iRow: " + str(iRow))
    # iRow is 0 indexed.
    if iRow == nRowsPerSlide:
      iRow = 0
    i = i + 1
  print ("Number of companies: ", i)
  powerpoint.save (finalPowerPoint)

  # Open the presentation file.
  # From https://stackoverflow.com/questions/17317219/is-there-an-platform-independent-equivalent-of-os-startfile/17317468#17317468

  open_file (finalPowerPoint)
except IOError:
  print(f"{bcolors.FAIL}IO Error (Wrong filename?): {bcolors.ENDC}" + inputFilename)
except IndexError as error:
  print(f"{bcolors.FAIL}Index Error (Forgot filename?):{bcolors.ENDC}", error)
except Exception as error:
  print(f"{bcolors.FAIL}Error:{bcolors.ENDC}", error)
